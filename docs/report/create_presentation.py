"""資源回収ルート最適化システムのPowerPointプレゼンテーション作成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """プレゼンテーションを作成"""
    prs = Presentation()
    prs.slide_width = Inches(16)  # ワイドスクリーン
    prs.slide_height = Inches(9)

    # スライド1: タイトル
    create_title_slide(prs)

    # スライド2: システム概要
    create_overview_slide(prs)

    # スライド3: システムフロー（簡潔版）
    create_flow_simple_slide(prs)

    # スライド4: 入力データ詳細
    create_input_slide(prs)

    # スライド5: 処理フロー詳細
    create_process_slide(prs)

    # スライド6: 出力結果詳細
    create_output_slide(prs)

    # スライド7: 技術スタックと特徴
    create_tech_slide(prs)

    # スライド8: まとめ
    create_summary_slide(prs)

    # 保存
    prs.save('claudedocs/system_presentation.pptx')
    print("PowerPoint file created: claudedocs/system_presentation.pptx")

def add_box(slide, left, top, width, height, text, fill_color, font_size=14, bold=True):
    """テキストボックスを追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    shape.line.color.rgb = RGBColor(100, 100, 100)
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.alignment = PP_ALIGN.CENTER

    return shape

def add_arrow(slide, x1, y1, x2, y2, color=(100, 100, 100)):
    """矢印を追加"""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        x1, y1, x2, y2
    )
    connector.line.color.rgb = RGBColor(*color)
    connector.line.width = Pt(3)
    return connector

def create_title_slide(prs):
    """タイトルスライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3),
        Inches(15), Inches(2)
    )
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "資源回収ルート最適化システム"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(21, 101, 192)
    p.alignment = PP_ALIGN.CENTER

    # サブタイトル
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5),
        Inches(15), Inches(1)
    )
    text_frame = subtitle_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "システムフロー図と技術概要"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # 日付
    date_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.5),
        Inches(15), Inches(0.5)
    )
    text_frame = date_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "2025年11月"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

def create_overview_slide(prs):
    """システム概要スライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "📋 システム概要"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(21, 101, 192)

    # 概要テキスト
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(6))
    tf = content_box.text_frame
    tf.word_wrap = True

    sections = [
        ("🎯 目的", [
            "資源回収業務におけるコスト最小化と効率化",
            "最短ルートで全回収地点を巡回",
            "エネルギー消費量（CO2排出量）の削減",
            "視覚的なルート表示と詳細なコスト分析"
        ]),
        ("💡 主な機能", [
            "地図クリックによる直感的な地点選択",
            "資源と車種の適合性自動チェック",
            "複数車両の同時最適化",
            "詳細なコスト内訳とエネルギー消費量表示"
        ]),
        ("👥 想定ユーザー", [
            "自治体の資源回収担当者",
            "廃棄物処理業者の配送計画担当者",
            "環境コンサルタント"
        ])
    ]

    y_offset = 0
    for title, items in sections:
        p = tf.paragraphs[0] if y_offset == 0 else tf.add_paragraph()
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(245, 124, 0)
        p.space_after = Pt(8)

        for item in items:
            p = tf.add_paragraph()
            p.text = f"  • {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(66, 66, 66)
            p.space_after = Pt(6)
            p.level = 1

        tf.add_paragraph()  # 空行

def create_flow_simple_slide(prs):
    """システムフロー簡潔版スライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "🔄 システムフロー概要"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(21, 101, 192)

    # 入力ボックス
    add_box(slide, Inches(0.8), Inches(1.5), Inches(3.5), Inches(5.5),
            "📥 入力\n\n道路ネットワーク\nデータ\n\nマスタデータ\n（資源・車種）\n\nユーザー選択\n（地点・量）",
            (227, 242, 253), 16, True)

    # 処理ボックス
    process_y = 1.8
    process_height = 0.7
    process_width = 4
    process_x = Inches(5.5)

    processes = [
        "1. データ読込・初期化",
        "2. 地点選択（UI）",
        "3. 車種割当",
        "4. 距離行列計算",
        "5. ルート最適化",
        "6. 結果生成"
    ]

    for i, proc in enumerate(processes):
        add_box(slide, process_x, Inches(process_y + i * 0.85),
                Inches(process_width), Inches(process_height),
                proc, (255, 243, 224), 14, True)

        # 矢印（次のステップへ）
        if i < len(processes) - 1:
            add_arrow(slide,
                     process_x + Inches(process_width/2),
                     Inches(process_y + i * 0.85 + process_height),
                     process_x + Inches(process_width/2),
                     Inches(process_y + (i+1) * 0.85),
                     (245, 124, 0))

    # 出力ボックス
    add_box(slide, Inches(11.7), Inches(1.5), Inches(3.5), Inches(5.5),
            "📤 出力\n\n最適ルート情報\n\nコスト詳細\n（固定費・変動費）\n\nエネルギー消費量\n\n地図表示",
            (232, 245, 233), 16, True)

    # 矢印（入力→処理）
    add_arrow(slide, Inches(4.3), Inches(4), Inches(5.5), Inches(4), (100, 100, 100))

    # 矢印（処理→出力）
    add_arrow(slide, Inches(9.5), Inches(4), Inches(11.7), Inches(4), (100, 100, 100))

def create_input_slide(prs):
    """入力データ詳細スライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "📥 入力データ詳細"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(21, 101, 192)

    # 3つの入力カテゴリ
    inputs = [
        {
            "title": "道路ネットワーク",
            "items": [
                "ファイル: road_network_*.json",
                "ノード情報（地点ID、緯度経度）",
                "エッジ情報（道路接続、距離）",
                "メタデータ"
            ],
            "x": 0.8,
            "color": (227, 242, 253)
        },
        {
            "title": "マスタデータ",
            "items": [
                "resources.json: 資源種別",
                "vehicles.json: 車種情報",
                "compatibility.json: 適合性",
                "嵩密度、コスト、制約条件"
            ],
            "x": 5.8,
            "color": (243, 229, 245)
        },
        {
            "title": "ユーザー選択",
            "items": [
                "車庫地点（出発・帰着）",
                "回収地点（複数地点）",
                "各地点の資源種別",
                "各地点の回収量[kg]",
                "集積場所（終点）"
            ],
            "x": 10.8,
            "color": (255, 243, 224)
        }
    ]

    for inp in inputs:
        # ボックス
        box = add_box(slide, Inches(inp["x"]), Inches(1.8),
                     Inches(4.2), Inches(5),
                     "", inp["color"], 14, True)

        # タイトル
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = inp["title"]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)

        # 項目
        for item in inp["items"]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(66, 66, 66)
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)

def create_process_slide(prs):
    """処理フロー詳細スライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "⚙️ 処理フロー詳細"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(245, 124, 0)

    # 処理ステップ
    processes = [
        ("1. データ読込", "JSON読込、グラフ構築\nキャッシュ初期化"),
        ("2. 地点選択", "地図クリック、最寄り検索\n資源種別・量入力"),
        ("3. 車種割当", "適合性チェック\n最適車種選択"),
        ("4. 距離計算", "最短経路探索\n距離行列生成"),
        ("5. 最適化", "VRP求解\nコスト最小化"),
        ("6. 結果生成", "経路再構成\n地図・表示")
    ]

    x_positions = [1, 6, 11, 1, 6, 11]
    y_positions = [1.5, 1.5, 1.5, 4.5, 4.5, 4.5]

    for i, (title, desc) in enumerate(processes):
        box = add_box(slide, Inches(x_positions[i]), Inches(y_positions[i]),
                     Inches(4), Inches(2),
                     "", (255, 243, 224), 14, True)

        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(230, 81, 0)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(66, 66, 66)
        p.alignment = PP_ALIGN.CENTER

        # 矢印
        if i < 2:  # 1→2, 2→3
            add_arrow(slide,
                     Inches(x_positions[i] + 4),
                     Inches(y_positions[i] + 1),
                     Inches(x_positions[i+1]),
                     Inches(y_positions[i+1] + 1),
                     (245, 124, 0))
        elif i == 2:  # 3→4（下へ）
            add_arrow(slide,
                     Inches(x_positions[i] + 2),
                     Inches(y_positions[i] + 2),
                     Inches(x_positions[i+1] + 2),
                     Inches(y_positions[i+1]),
                     (245, 124, 0))
        elif i < 5:  # 4→5, 5→6
            add_arrow(slide,
                     Inches(x_positions[i] + 4),
                     Inches(y_positions[i] + 1),
                     Inches(x_positions[i+1]),
                     Inches(y_positions[i+1] + 1),
                     (245, 124, 0))

def create_output_slide(prs):
    """出力結果詳細スライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "📤 出力結果詳細"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 125, 50)

    # 4つの出力カテゴリ
    outputs = [
        {
            "title": "最適ルート情報",
            "items": [
                "訪問順序リスト",
                "総走行距離 [km]",
                "使用車種名",
                "各区間の距離"
            ],
            "x": 0.8,
            "y": 1.8
        },
        {
            "title": "コスト詳細内訳",
            "items": [
                "固定費（項目別）",
                "  - 人件費、車両償却費等",
                "変動費（項目別）",
                "  - 燃料費、修繕費等",
                "総コスト [円]"
            ],
            "x": 8.4,
            "y": 1.8
        },
        {
            "title": "エネルギー消費量",
            "items": [
                "総消費電力量 [kWh]",
                "車両別消費量",
                "CO2削減効果",
                "（EV化の場合）"
            ],
            "x": 0.8,
            "y": 4.8
        },
        {
            "title": "地図上の可視化",
            "items": [
                "インタラクティブ地図",
                "ルート経路表示（青線）",
                "地点マーカー（色分け）",
                "訪問順序番号",
                "凡例表示"
            ],
            "x": 8.4,
            "y": 4.8
        }
    ]

    for out in outputs:
        box = add_box(slide, Inches(out["x"]), Inches(out["y"]),
                     Inches(6.8), Inches(2.5),
                     "", (232, 245, 233), 14, True)

        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = out["title"]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(27, 94, 32)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(10)

        for item in out["items"]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(66, 66, 66)
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(5)

def create_tech_slide(prs):
    """技術スタックと特徴スライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "🔧 技術スタックとシステム特徴"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(21, 101, 192)

    # 左側：技術スタック
    tech_box = add_box(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(6),
                      "", (255, 248, 225), 14, True)
    tf = tech_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "技術スタック"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(245, 124, 0)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(15)

    tech_items = [
        ("Streamlit", "Webアプリケーションフレームワーク"),
        ("NetworkX", "グラフ処理・最短経路計算"),
        ("Folium", "インタラクティブ地図表示"),
        ("Pandas", "データ処理・テーブル編集"),
        ("Python 3.8+", "システム全体の実装言語")
    ]

    for tech, desc in tech_items:
        p = tf.add_paragraph()
        p.text = f"● {tech}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(3)

        p = tf.add_paragraph()
        p.text = f"   {desc}"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.space_after = Pt(10)

    # 右側：システム特徴
    feature_box = add_box(slide, Inches(8.2), Inches(1.5), Inches(7), Inches(6),
                         "", (227, 242, 253), 14, True)
    tf = feature_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "システム特徴"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(21, 101, 192)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(15)

    features = [
        "地図クリックで直感的な地点選択",
        "資源と車種の適合性を自動チェック",
        "複数車両の同時最適化に対応",
        "詳細なコスト内訳表示",
        "エネルギー消費量の可視化",
        "インタラクティブ地図でルート表示",
        "セッション状態保持で効率的操作"
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = f"✓ {feature}"
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(8)

def create_summary_slide(prs):
    """まとめスライドを作成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "📌 まとめ"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(21, 101, 192)

    # まとめ内容
    summary_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(13), Inches(5))
    tf = summary_box.text_frame
    tf.word_wrap = True

    sections = [
        ("🎯 システムの価値", [
            "道路ネットワーク上での資源回収ルートを最適化",
            "コストとエネルギー消費を最小化し、環境負荷を低減",
            "直感的なUIで専門知識がなくても利用可能"
        ]),
        ("📊 主要な処理", [
            "入力: 道路ネットワーク、マスタデータ、ユーザー選択",
            "処理: 距離計算 → VRP最適化 → コスト算出",
            "出力: 最適ルート、詳細コスト、地図表示"
        ]),
        ("🚀 今後の展開", [
            "複数日のスケジューリング対応",
            "リアルタイム交通情報の統合",
            "機械学習による需要予測機能"
        ])
    ]

    for title, items in sections:
        p = tf.paragraphs[0] if title == sections[0][0] else tf.add_paragraph()
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(245, 124, 0)
        p.space_after = Pt(10)

        for item in items:
            p = tf.add_paragraph()
            p.text = f"  • {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(66, 66, 66)
            p.space_after = Pt(8)
            p.level = 1

        tf.add_paragraph()

if __name__ == "__main__":
    create_presentation()
