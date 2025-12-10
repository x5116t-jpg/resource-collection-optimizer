"""SVGフロー図を1枚のPowerPointスライドに変換"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_single_slide_presentation():
    """1枚のスライドにフロー図を作成"""
    prs = Presentation()
    prs.slide_width = Inches(16)  # ワイドスクリーン
    prs.slide_height = Inches(9)

    # 空白スライドを追加
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ========== タイトル ==========
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "資源回収ルート最適化システム"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = RGBColor(21, 101, 192)
    p.alignment = PP_ALIGN.CENTER

    # サブタイトル
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(15), Inches(0.4))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "システムフロー概要図"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER

    # ========== システム境界ボックス ==========
    boundary = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(1.3),
        Inches(15.2), Inches(6.8)
    )
    boundary.fill.solid()
    boundary.fill.fore_color.rgb = RGBColor(250, 250, 250)
    boundary.line.color.rgb = RGBColor(25, 118, 210)
    boundary.line.width = Pt(3)
    boundary.line.dash_style = 3  # Dash

    # ========== 入力ボックス ==========
    input_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.6),
        Inches(3.5), Inches(2.2)
    )
    input_box.fill.solid()
    input_box.fill.fore_color.rgb = RGBColor(227, 242, 253)
    input_box.line.color.rgb = RGBColor(25, 118, 210)
    input_box.line.width = Pt(3)

    tf = input_box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = "📥 入力"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(13, 71, 161)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)

    items = [
        "● 道路ネットワークデータ",
        "● マスタデータ",
        "● ユーザー選択"
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(21, 101, 192)
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(5)

    # ========== 処理ボックス群 ==========
    process_x = 4.8
    process_y_start = 1.6
    process_height = 0.65
    process_width = 3.0
    process_gap = 0.08

    processes = [
        "1. データ読込・初期化",
        "2. 地点選択（UI操作）",
        "3. 車種割当プラン生成",
        "4. 距離行列計算",
        "5. ルート最適化（VRP）",
        "6. 結果生成・可視化"
    ]

    for i, proc_text in enumerate(processes):
        proc_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(process_x), Inches(process_y_start + i * (process_height + process_gap)),
            Inches(process_width), Inches(process_height)
        )
        proc_box.fill.solid()
        proc_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        proc_box.line.color.rgb = RGBColor(245, 124, 0)
        proc_box.line.width = Pt(2)

        tf = proc_box.text_frame
        tf.vertical_anchor = 1  # Middle
        p = tf.paragraphs[0]
        p.text = proc_text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(230, 81, 0)
        p.alignment = PP_ALIGN.CENTER

        # 矢印（次のステップへ）
        if i < len(processes) - 1:
            arrow = slide.shapes.add_connector(
                1,  # STRAIGHT
                Inches(process_x + process_width/2),
                Inches(process_y_start + i * (process_height + process_gap) + process_height),
                Inches(process_x + process_width/2),
                Inches(process_y_start + (i+1) * (process_height + process_gap))
            )
            arrow.line.color.rgb = RGBColor(245, 124, 0)
            arrow.line.width = Pt(3)

    # 処理全体のラベル
    proc_label = slide.shapes.add_textbox(
        Inches(process_x), Inches(1.35),
        Inches(process_width), Inches(0.25)
    )
    tf = proc_label.text_frame
    p = tf.paragraphs[0]
    p.text = "⚙️ 処理"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 81, 0)
    p.alignment = PP_ALIGN.CENTER

    # ========== 出力ボックス ==========
    output_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.3), Inches(1.6),
        Inches(3.5), Inches(2.2)
    )
    output_box.fill.solid()
    output_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    output_box.line.color.rgb = RGBColor(56, 142, 60)
    output_box.line.width = Pt(3)

    tf = output_box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = "📤 出力"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(27, 94, 32)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)

    items = [
        "● 最適ルート情報",
        "● コスト詳細",
        "● エネルギー消費量",
        "● 地図表示・分析結果"
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(46, 125, 50)
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(5)

    # ========== 矢印（入力→処理） ==========
    arrow1 = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(4.3), Inches(2.7),
        Inches(4.8), Inches(2.7)
    )
    arrow1.line.color.rgb = RGBColor(100, 100, 100)
    arrow1.line.width = Pt(4)

    # ========== 矢印（処理→出力） ==========
    arrow2 = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(7.8), Inches(2.7),
        Inches(8.3), Inches(2.7)
    )
    arrow2.line.color.rgb = RGBColor(100, 100, 100)
    arrow2.line.width = Pt(4)

    # ========== データ格納層 ==========
    data_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(4.2),
        Inches(3.5), Inches(2.5)
    )
    data_box.fill.solid()
    data_box.fill.fore_color.rgb = RGBColor(243, 229, 245)
    data_box.line.color.rgb = RGBColor(123, 31, 162)
    data_box.line.width = Pt(2)

    tf = data_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "💾 データ格納層"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(74, 20, 140)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)

    items = [
        "● NetworkXグラフ",
        "● 地点レジストリ",
        "● 車種カタログ",
        "● 距離行列",
        "● 空間インデックス"
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(106, 27, 154)
        p.space_after = Pt(3)

    # ========== 技術情報 ==========
    tech_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.3), Inches(4.2),
        Inches(3.5), Inches(2.5)
    )
    tech_box.fill.solid()
    tech_box.fill.fore_color.rgb = RGBColor(255, 248, 225)
    tech_box.line.color.rgb = RGBColor(245, 124, 0)
    tech_box.line.width = Pt(2)

    tf = tech_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "🔧 技術スタック"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 81, 0)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)

    items = [
        "● Streamlit（Webアプリ）",
        "● NetworkX（グラフ処理）",
        "● Folium（地図表示）",
        "● Pandas（データ処理）"
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(239, 108, 0)
        p.space_after = Pt(3)

    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(5)

    p = tf.add_paragraph()
    p.text = "最適化: VRP（配送計画問題）"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(216, 67, 21)

    # ========== システム特徴（下部） ==========
    feature_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(7),
        Inches(11), Inches(0.8)
    )
    feature_box.fill.solid()
    feature_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    feature_box.line.color.rgb = RGBColor(153, 153, 153)
    feature_box.line.width = Pt(1)

    tf = feature_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "システムの特徴"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)

    features_text = "✓ 地図クリックで直感的に地点選択   ✓ 資源と車種の適合性を自動チェック   ✓ 複数車両の同時最適化   ✓ 詳細なコスト・エネルギー分析   ✓ インタラクティブ地図でルート可視化"

    p = tf.add_paragraph()
    p.text = features_text
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.alignment = PP_ALIGN.CENTER

    # ========== フッター ==========
    footer = slide.shapes.add_textbox(
        Inches(0.5), Inches(8.2),
        Inches(15), Inches(0.3)
    )
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "本システムは道路ネットワーク上での資源回収ルートを最適化し、コストとエネルギー消費を最小化します"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(153, 153, 153)
    p.alignment = PP_ALIGN.CENTER

    # 保存
    prs.save('claudedocs/system_flowchart_slide.pptx')
    print("PowerPoint slide created: claudedocs/system_flowchart_slide.pptx")

if __name__ == "__main__":
    create_single_slide_presentation()
